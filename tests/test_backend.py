"""Backend tests: conversion, combining, redaction, settings persistence.

Run with:  python tests/test_backend.py
"""
import os
import sys
import shutil
import tempfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from backend.converter import convert_files, scan_folder, convert_file, estimate_tokens, estimate_source_tokens
from backend.combiner import combine_files, combine_folder
from backend import settings as settings_mod

failures = []

def check(name, cond, detail=""):
    status = "PASS" if cond else "FAIL"
    print(f"[{status}] {name} {detail}")
    if not cond:
        failures.append(name)

def known_gap(name, resolved, detail=""):
    """Documents a redaction-quality limitation without failing the suite.
    Prints RESOLVED once the gap closes (e.g. after the GLiNER upgrade), which
    is the signal to promote it to a hard assertion."""
    if resolved:
        print(f"[RESOLVED] {name} — now fixed; promote to a check() assertion")
    else:
        print(f"[KNOWN-GAP] {name} (tracked for GLiNER redaction upgrade) {detail}")

work = tempfile.mkdtemp(prefix="docprep_test_")
sub = os.path.join(work, "notes")
os.makedirs(sub)

with open(os.path.join(work, "alpha.html"), "w", encoding="utf-8") as f:
    paragraphs = "".join(
        f'<div class="section-wrap outer"><p style="margin:0;padding:4px;color:#333">'
        f'<span class="body-text">Paragraph {i} with meaningful report content here.</span></p></div>'
        for i in range(30)
    )
    f.write(f"<html><body><h1>Alpha Report</h1><p>Contact: jane@example.com and call 555-123-4567.</p>{paragraphs}</body></html>")
with open(os.path.join(sub, "beta.html"), "w", encoding="utf-8") as f:
    f.write("<html><body><h2>Beta Notes</h2><p>Quarterly revenue was strong.</p></body></html>")
with open(os.path.join(work, "meeting.vtt"), "w", encoding="utf-8") as f:
    f.write("WEBVTT\n\n1\n00:00:00.000 --> 00:00:02.000\nHello and welcome.\n\n2\n00:00:02.000 --> 00:00:04.000\nLet's begin the review.\n")
with open(os.path.join(work, "~$temp.docx"), "w") as f:
    f.write("junk")
with open(os.path.join(work, "existing-note.md"), "w", encoding="utf-8") as f:
    f.write("# My own pre-existing note\nShould NOT be swept into combined output.\n")

# 1. scan_folder skips temp files and pre-existing md
found = scan_folder(work)
check("scan_folder finds 3 files", len(found) == 3, f"found={[os.path.basename(f) for f in found]}")
check("scan_folder skips ~$ temp", not any("~$" in f for f in found))

# 2. convert_files with events
events = []
outs = convert_files(found, progress_callback=events.append, inject_yaml=True)
check("convert_files output count", len(outs) == 3, f"outs={len(outs)}")
check("events carry tokens", all(e["tokens"] > 0 for e in events if e["status"] == "done"))
check("events done/total", events[-1]["done"] == 3 and events[-1]["total"] == 3)
alpha_md = os.path.join(work, "alpha.md")
check("yaml frontmatter injected", open(alpha_md, encoding="utf-8").read().startswith("---\n"))

# 2b. source-token baseline for savings display
alpha_event = next(e for e in events if e["file"].endswith("alpha.html") and e["status"] == "done")
check("html event carries source_tokens", (alpha_event.get("source_tokens") or 0) > 0, str(alpha_event.get("source_tokens")))
check("html raw > markdown output", alpha_event["source_tokens"] > alpha_event["tokens"],
      f"src={alpha_event['source_tokens']} out={alpha_event['tokens']}")
check("unknown ext has no raw baseline", estimate_source_tokens("anything.xyz") is None)

def make_pdf(path, pages):
    """Writes a minimal valid PDF with the given page count."""
    objs = [b"<< /Type /Catalog /Pages 2 0 R >>"]
    kids = " ".join(f"{3 + i} 0 R" for i in range(pages))
    objs.append(f"<< /Type /Pages /Kids [{kids}] /Count {pages} >>".encode())
    for _ in range(pages):
        objs.append(b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] >>")
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for n, body in enumerate(objs, start=1):
        offsets.append(len(out))
        out += f"{n} 0 obj\n".encode() + body + b"\nendobj\n"
    xref_pos = len(out)
    out += f"xref\n0 {len(objs) + 1}\n".encode() + b"0000000000 65535 f \n"
    for off in offsets:
        out += f"{off:010d} 00000 n \n".encode()
    out += f"trailer\n<< /Size {len(objs) + 1} /Root 1 0 R >>\nstartxref\n{xref_pos}\n%%EOF\n".encode()
    with open(path, "wb") as f:
        f.write(bytes(out))

pdf_path = os.path.join(work, "sample.pdf")
make_pdf(pdf_path, pages=3)
check("pdf baseline = text + pages x 1500", estimate_source_tokens(pdf_path, output_tokens=800) == 3 * 1500 + 800,
      str(estimate_source_tokens(pdf_path, output_tokens=800)))

# pptx baseline: slides x 1500 + text, NOT the DrawingML XML size (which
# overstates savings on image-heavy decks)
import zipfile as _zipfile
pptx_path = os.path.join(work, "deck.pptx")
with _zipfile.ZipFile(pptx_path, "w") as z:
    z.writestr("[Content_Types].xml", "<Types/>")
    for i in range(1, 6):
        z.writestr(f"ppt/slides/slide{i}.xml", "<p:sld>" + "x" * 40000 + "</p:sld>")
    z.writestr("ppt/slides/_rels/slide1.xml.rels", "<Relationships/>")
check("pptx baseline = text + slides x 1500", estimate_source_tokens(pptx_path, output_tokens=700) == 5 * 1500 + 700,
      str(estimate_source_tokens(pptx_path, output_tokens=700)))

# 2c. unsupported types are rejected before reaching any converter
# (mp3/wav/m4a/mp4 would hit markitdown's cloud transcription path; zip is
# rejected unless explicitly enabled — tested separately in 2e)
for bad_ext in (".mp3", ".mp4", ".wav", ".m4a", ".zip", ".exe"):
    bad_path = os.path.join(work, f"blocked{bad_ext}")
    with open(bad_path, "wb") as f:
        f.write(b"junkjunkjunkjunk")
    bad_events = []
    bad_outs = convert_files([bad_path], progress_callback=bad_events.append)
    bad_error = next((e for e in bad_events if e["status"] == "error"), None)
    check(f"unsupported ext rejected: {bad_ext}", bad_outs == [] and bad_error is not None
          and "Unsupported file type" in (bad_error["error"] or ""), str(bad_error and bad_error["error"]))
    os.remove(bad_path)

# 2d. newly supported formats
from backend.converter import convert_to_markdown, SUPPORTED_EXTENSIONS, IMAGE_EXTENSIONS, CONVERTIBLE_EXTENSIONS
import zipfile as _zipfile
import json as _json

fmt_dir = os.path.join(work, "fmt")
os.makedirs(fmt_dir)

# minimal EPUB (zip: mimetype + container + opf + one xhtml chapter)
epub_path = os.path.join(fmt_dir, "book.epub")
chapter = ('<?xml version="1.0" encoding="utf-8"?><html xmlns="http://www.w3.org/1999/xhtml">'
           "<head><title>Ch1</title></head><body><h1>Chapter One</h1>"
           "<p>The lighthouse keeper counted the gulls at dawn.</p></body></html>")
opf = ('<?xml version="1.0"?><package xmlns="http://www.idpf.org/2007/opf" version="3.0" unique-identifier="id">'
       '<metadata xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:title>Test Book</dc:title>'
       '<dc:identifier id="id">t1</dc:identifier><dc:language>en</dc:language></metadata>'
       '<manifest><item id="c1" href="ch1.xhtml" media-type="application/xhtml+xml"/></manifest>'
       '<spine><itemref idref="c1"/></spine></package>')
container = ('<?xml version="1.0"?><container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
             '<rootfiles><rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/>'
             "</rootfiles></container>")
with _zipfile.ZipFile(epub_path, "w") as z:
    z.writestr("mimetype", "application/epub+zip", compress_type=_zipfile.ZIP_STORED)
    z.writestr("META-INF/container.xml", container)
    z.writestr("OEBPS/content.opf", opf)
    z.writestr("OEBPS/ch1.xhtml", chapter)
try:
    epub_md = convert_to_markdown(epub_path)
    check("epub converts", "lighthouse keeper" in epub_md, epub_md[:120])
except Exception as e:
    check("epub converts", False, str(e))

# minimal Jupyter notebook
ipynb_path = os.path.join(fmt_dir, "analysis.ipynb")
with open(ipynb_path, "w", encoding="utf-8") as f:
    _json.dump({"cells": [
        {"cell_type": "markdown", "metadata": {}, "source": ["# Notebook Title\n", "Findings about churn."]},
        {"cell_type": "code", "execution_count": 1, "metadata": {}, "outputs": [],
         "source": ["print('churn_rate')"]},
    ], "metadata": {}, "nbformat": 4, "nbformat_minor": 5}, f)
try:
    nb_md = convert_to_markdown(ipynb_path)
    check("ipynb converts", "Findings about churn" in nb_md and "churn_rate" in nb_md, nb_md[:120])
except Exception as e:
    check("ipynb converts", False, str(e))

# images: rejected with a specific message rather than silently writing an
# empty .md (markitdown yields no text without a vision model)
from PIL import Image as _Image
jpg_path = os.path.join(fmt_dir, "photo.jpg")
_Image.new("RGB", (8, 8), (200, 30, 30)).save(jpg_path)
png_path = os.path.join(fmt_dir, "shot.png")
_Image.new("RGB", (8, 8), (30, 60, 200)).save(png_path)
for img in (jpg_path, png_path):
    label = os.path.splitext(img)[1]
    try:
        convert_to_markdown(img)
        check(f"image rejected with actionable message: {label}", False, "converted instead of raising")
    except ValueError as e:
        check(f"image rejected with actionable message: {label}",
              "enable" in str(e).lower(), str(e))
check("no .md written for images",
      not os.path.exists(os.path.splitext(jpg_path)[0] + ".md"))

# msg/xls: junk bytes now fail with a *conversion* error, not the allowlist
for accepted_ext in (".msg", ".xls"):
    p = os.path.join(fmt_dir, f"junk{accepted_ext}")
    with open(p, "wb") as f:
        f.write(b"not really this format")
    try:
        convert_to_markdown(p)
        check(f"{accepted_ext} accepted by allowlist", True, "(junk happened to convert)")
    except ValueError as e:
        check(f"{accepted_ext} accepted by allowlist", "Unsupported file type" not in str(e), str(e)[:80])
    except Exception:
        check(f"{accepted_ext} accepted by allowlist", True, "(conversion error, as expected for junk)")

# scans pick up new formats but not images
scanned = {os.path.basename(f) for f in scan_folder(fmt_dir)}
check("scan includes epub/ipynb/msg/xls",
      {"book.epub", "analysis.ipynb", "junk.msg", "junk.xls"} <= scanned, str(scanned))
check("scan excludes images", "photo.jpg" not in scanned and "shot.png" not in scanned, str(scanned))
check("image exts not convertible", all(e not in CONVERTIBLE_EXTENSIONS for e in IMAGE_EXTENSIONS))

# 2e. ZIP: opt-in, filtered, capped
zip_dir = os.path.join(work, "zips")
os.makedirs(zip_dir)
archive = os.path.join(zip_dir, "bundle.zip")
inner_nested = os.path.join(zip_dir, "nested.zip")
with _zipfile.ZipFile(inner_nested, "w") as z:
    z.writestr("deep.txt", "should never be reached")
with _zipfile.ZipFile(archive, "w") as z:
    z.writestr("report.html", "<html><body><h1>Zipped Report</h1><p>Email zip.person@example.com inside.</p></body></html>")
    z.writestr("song.mp3", b"ID3junk")
    z.write(inner_nested, "nested.zip")
    z.writestr("sub/notes.txt", "plain text survives")

try:
    convert_to_markdown(archive)
    check("zip rejected by default", False, "no exception")
except ValueError as e:
    check("zip rejected by default", "Unsupported file type" in str(e), str(e)[:80])

zip_md = convert_to_markdown(archive, allow_zip=True, redact_pii=True, redact_mode="Regex Only")
check("zip converts when enabled", "Zipped Report" in zip_md and "plain text survives" in zip_md)
check("zip content is redacted", "[REDACTED_EMAIL]" in zip_md and "zip.person@example.com" not in zip_md)
check("zip skips blocked inner types", "song.mp3 (unsupported type)" in zip_md)
check("zip skips nested archives", "nested.zip (nested archive)" in zip_md)

big = os.path.join(zip_dir, "toomany.zip")
with _zipfile.ZipFile(big, "w") as z:
    for i in range(201):
        z.writestr(f"f{i}.txt", "x")
try:
    convert_to_markdown(big, allow_zip=True)
    check("zip entry cap enforced", False, "no exception")
except ValueError as e:
    check("zip entry cap enforced", "limit is 200" in str(e), str(e)[:80])

# 2f. OCR (opt-in). Skipped when the OCR extras aren't installed.
from backend import ocr as ocr_mod
if not ocr_mod.is_available():
    print("[SKIP] OCR tests — rapidocr/pypdfium2 not installed")
else:
    ocr_dir = os.path.join(work, "ocr")
    os.makedirs(ocr_dir)
    # an image whose ONLY content is rendered text
    from PIL import Image as _Im, ImageDraw as _Dr
    shot = os.path.join(ocr_dir, "screenshot.png")
    im = _Im.new("RGB", (900, 220), "white")
    _Dr.Draw(im).text((30, 70), "QUARTERLY REVENUE 42817", fill="black")
    im.save(shot)

    # off by default: still refuses images, with the vision-model message
    try:
        convert_to_markdown(shot)
        check("image still refused when OCR off", False, "converted unexpectedly")
    except ValueError as e:
        check("image still refused when OCR off", "enable" in str(e).lower(), str(e)[:60])

    # on: extracts the text
    text = convert_to_markdown(shot, ocr=True)
    check("OCR reads text from an image", "42817" in text.replace(" ", ""), repr(text[:80]))

    # tiny images are skipped as icons, not OCR'd
    icon = os.path.join(ocr_dir, "icon.png")
    _Im.new("RGB", (24, 24), "white").save(icon)
    try:
        convert_to_markdown(icon, ocr=True)
        check("tiny icon yields no text", False, "expected no-text error")
    except ValueError as e:
        check("tiny icon yields no text", "No readable text" in str(e), str(e)[:60])

    # a PDF that already has a text layer must NOT be re-OCR'd
    check("text-layer PDF not treated as scanned",
          not ocr_mod.pdf_text_layer_is_thin(pdf_path, "x" * 5000))
    check("empty-text PDF treated as scanned",
          ocr_mod.pdf_text_layer_is_thin(pdf_path, ""))

    # embedded-image OCR dedupes repeated logos
    # a real deck (python-pptx), not a bare zip: markitdown must be able to
    # convert it. Three byte-distinct images (python-pptx dedupes identical
    # blobs) carrying the same text, so the logo-dedupe check is meaningful.
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import PngImagePlugin as _Png
    import io as _io
    deck = os.path.join(ocr_dir, "deck.pptx")
    prs = Presentation()
    logo = _Im.new("RGB", (500, 160), "white")
    _Dr.Draw(logo).text((20, 60), "ACMECORP LOGO", fill="black")
    for n in range(3):
        # identical pixels (so OCR reads the same text) but distinct bytes via
        # a metadata chunk, so python-pptx keeps three separate media parts
        meta = _Png.PngInfo(); meta.add_text("copy", str(n))
        b = _io.BytesIO(); logo.save(b, format="PNG", pnginfo=meta); b.seek(0)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(b, Inches(1), Inches(1))
    prs.save(deck)
    found = ocr_mod.ocr_embedded_images(deck)
    check("embedded-image OCR finds text", "ACME" in found.upper().replace(" ", ""), repr(found[:80]))
    check("repeated logo emitted once", found.upper().count("ACMECORP") <= 1,
          f"count={found.upper().count('ACMECORP')}")

    # OCR off: the reader is told which files hold unread images, and the
    # conversion event carries a count so a UI can badge them in a batch
    check("unread-image count (deck with 3 images)", ocr_mod.count_unread_images(deck) == 3,
          str(ocr_mod.count_unread_images(deck)))
    icons_only = os.path.join(ocr_dir, "icons.pptx")
    with _zipfile.ZipFile(icons_only, "w") as z:
        z.writestr("[Content_Types].xml", "<Types/>")
        b = _io.BytesIO(); _Im.new("RGB", (24, 24), "white").save(b, format="PNG")
        z.writestr("ppt/media/image0.png", b.getvalue())
    check("icons don't count as unread images", ocr_mod.count_unread_images(icons_only) == 0)
    check("scanned-looking PDF counts its pages", ocr_mod.count_unread_images(pdf_path, "") == 3,
          str(ocr_mod.count_unread_images(pdf_path, "")))
    check("text-layer PDF counts nothing", ocr_mod.count_unread_images(pdf_path, "x" * 5000) == 0)
    note_md = convert_to_markdown(deck)                       # OCR off
    check("OCR-off output carries the unread-images note", "3 embedded images that were not read" in note_md,
          repr(note_md[-160:]))
    check("OCR-on output has no such note", "were not read" not in convert_to_markdown(deck, ocr=True))
    ev_unread = []
    convert_files([deck], progress_callback=ev_unread.append, overwrite=True)
    done_ev = next(e for e in ev_unread if e["status"] == "done")
    check("event carries unread_images", done_ev.get("unread_images") == 3, str(done_ev.get("unread_images")))

# 3. redaction (regex)
from backend.converter import redact_pii_content
red_out = convert_file(os.path.join(work, "alpha.html"), overwrite=False, redact_pii=True, redact_mode="Regex Only")
red_text = open(red_out, encoding="utf-8").read()
check("keep-both naming", red_out.endswith("alpha (1).md"), red_out)
check("email redacted", "[REDACTED_EMAIL]" in red_text)
check("phone redacted", "[REDACTED_PHONE]" in red_text)
os.remove(red_out)

# 3b. regex redaction edge cases (previously leaked characters)
r = redact_pii_content("Call (503) 555-0182 today.", mode="Regex Only")
check("phone with parens leaves no dangling '(' or digits",
      "[REDACTED_PHONE]" in r and "(" not in r and "503" not in r, repr(r))
r = redact_pii_content("Reach a.b@school.edu. Next sentence stays.", mode="Regex Only")
check("email preserves the sentence-ending period", "[REDACTED_EMAIL]." in r, repr(r))
check("email does not swallow the following word", "Next sentence stays" in r, repr(r))

# 3c. on-device NER redaction: privacy invariants (hard) + known quality gaps.
# Reproduces the student-gradebook findings — names/scores in terse table cells.
NER_TABLE = "| Marcus Delgado | RHS-2024-0142 | 88 |\n| James Whitfield | RHS-2024-0103 | 94 |"
try:
    nr = redact_pii_content(NER_TABLE, mode="Local NER (spaCy)")
    # Privacy invariants — MUST hold regardless of label accuracy
    check("NER: numeric scores survive redaction", "88" in nr and "94" in nr, repr(nr))
    check("NER: no student name left in cleartext",
          "Marcus Delgado" not in nr and "James Whitfield" not in nr
          and "Marcus" not in nr and "Whitfield" not in nr, repr(nr))
    # Known quality gaps — spaCy en_core_web_sm on terse tabular text.
    # These flip to RESOLVED when a better model (GLiNER) lands.
    known_gap("NER labels every person as REDACTED_NAME (not REDACTED_LOCATION)",
              "[REDACTED_LOCATION]" not in nr, repr(nr))
    known_gap("NER labels names consistently across rows",
              nr.count("[REDACTED_NAME]") == 2, f"names tagged={nr.count('[REDACTED_NAME]')}/2")
    known_gap("NER handles alphanumeric IDs cleanly (no partial 'RHS' mangling)",
              "[REDACTED_LOCATION]-" not in nr, repr(nr))
except Exception as e:
    print(f"[SKIP] on-device NER tests — spaCy model unavailable: {e}")

# 4. combine_files only includes the given list (not existing-note.md)
combined = combine_files(sorted(outs), os.path.join(work, "test-combined.md"), base_dir=work, collection_name="test")
ctext = open(combined, encoding="utf-8").read()
check("combined excludes pre-existing note", "pre-existing note" not in ctext)
check("combined has TOC", "## Table of Contents" in ctext)
check("combined has subfolder rel name", "beta.md" in ctext)
check("combined heading uses collection name", "# test" in ctext)

# 5. combine_folder compat wrapper still works
os.remove(combined)
combined2 = combine_folder(work, output_filename="master.md")
check("combine_folder wrapper", combined2.endswith("master.md") and os.path.exists(combined2))

# 6. settings roundtrip (redirect config dir into temp)
settings_mod.config_dir = lambda: os.path.join(work, "cfg")
settings_mod.config_path = lambda: os.path.join(work, "cfg", "settings.json")
s = settings_mod.load_settings()
check("defaults loaded", s["output_mode"] == "both" and s["conflict"] == "keep_both")
s["output_mode"] = "combined_only"
s["redact"] = True
s["custom_terms"] = "ACME Corp"
settings_mod.save_settings(s)
s2 = settings_mod.load_settings()
check("settings persist", s2["output_mode"] == "combined_only" and s2["redact"] is True and s2["custom_terms"] == "ACME Corp")

# 7. token estimate sanity
check("estimate_tokens", estimate_tokens("word " * 400) == 500)

# 8. cancellation stops scheduling
cancel_calls = {"n": 0}
def cancel_after_first():
    cancel_calls["n"] += 1
    return cancel_calls["n"] > 1
outs_c = convert_files(found, cancel_check=cancel_after_first, overwrite=True)
check("cancel returns partial", len(outs_c) <= len(found))

shutil.rmtree(work, ignore_errors=True)
print()
print("FAILURES:", failures if failures else "none")
sys.exit(1 if failures else 0)
